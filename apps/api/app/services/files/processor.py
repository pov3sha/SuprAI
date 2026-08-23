import io
import json
import csv
from typing import List, Dict, Any
from sqlalchemy.orm import Session
from loguru import logger
from pypdf import PdfReader
from app.models.schema import FileRecord, FileChunk
from app.services.files.storage import storage_service

def process_uploaded_document(file_id: str, db: Session) -> FileRecord:
    file_record = db.query(FileRecord).filter(FileRecord.id == file_id).first()
    if not file_record:
        logger.error(f"FileRecord {file_id} not found")
        return None

    # Fetch object from Storage Service
    file_bytes = storage_service.get_file(file_record.storage_path)
    filename = file_record.filename.lower()

    chunks = []
    page_count = 1

    if filename.endswith(".pdf"):
        chunks, page_count = parse_pdf(file_id, file_bytes)
    elif filename.endswith(".docx"):
        chunks, page_count = parse_docx(file_id, file_bytes)
    elif filename.endswith(".pptx"):
        chunks, page_count = parse_pptx(file_id, file_bytes)
    elif filename.endswith(".csv"):
        chunks, page_count = parse_csv_dataset(file_id, file_bytes)
    elif filename.endswith(".xlsx") or filename.endswith(".xls"):
        chunks, page_count = parse_excel_dataset(file_id, file_bytes)
    elif filename.endswith(".json"):
        chunks, page_count = parse_json_document(file_id, file_bytes)
    elif filename.endswith(".txt") or filename.endswith(".md"):
        chunks, page_count = parse_text_document(file_id, file_bytes)
    else:
        logger.warning(f"Unsupported file format: {file_record.filename}")
        file_record.status = "UNSUPPORTED"
        db.commit()
        raise ValueError(f"Unsupported file type: {file_record.filename}")

    # Persist Chunks into PostgreSQL
    db.query(FileChunk).filter(FileChunk.file_id == file_id).delete()
    for chunk_data in chunks:
        chunk = FileChunk(
            file_id=file_id,
            page_number=chunk_data["source_location"],
            content=chunk_data["content"],
            metadata_json={
                "source_type": chunk_data["source_type"],
                "source_location": chunk_data["source_location"],
                "file_type": file_record.mime_type
            }
        )
        db.add(chunk)

    file_record.page_count = page_count
    file_record.status = "PROCESSED"
    db.commit()
    db.refresh(file_record)
    return file_record

def parse_pdf(file_id: str, file_bytes: bytes):
    reader = PdfReader(io.BytesIO(file_bytes))
    chunks = []
    for idx, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            chunks.append({
                "source_type": "page",
                "source_location": idx + 1,
                "content": text.strip()
            })
    return chunks, len(reader.pages)

def parse_docx(file_id: str, file_bytes: bytes):
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        chunks = []
        for idx, p in enumerate(doc.paragraphs):
            if p.text.strip():
                chunks.append({
                    "source_type": "paragraph",
                    "source_location": idx + 1,
                    "content": p.text.strip()
                })
        return chunks, len(chunks)
    except Exception as e:
        logger.warning(f"python-docx parse fallback: {e}")
        return parse_text_document(file_id, file_bytes)

def parse_pptx(file_id: str, file_bytes: bytes):
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        chunks = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            full_text = "\n".join(slide_text)
            if full_text.strip():
                chunks.append({
                    "source_type": "slide",
                    "source_location": idx + 1,
                    "content": full_text.strip()
                })
        return chunks, len(prs.slides)
    except Exception as e:
        logger.warning(f"python-pptx parse fallback: {e}")
        return parse_text_document(file_id, file_bytes)

def parse_csv_dataset(file_id: str, file_bytes: bytes):
    text_content = file_bytes.decode('utf-8', errors='ignore')
    reader = csv.reader(io.StringIO(text_content))
    rows = list(reader)
    if not rows:
        return [], 0

    header = rows[0]
    row_count = len(rows) - 1
    
    summary = (
        f"DATASET SUMMARY (CSV):\n"
        f"Total Rows: {row_count}\n"
        f"Columns ({len(header)}): {', '.join(header[:15])}\n\n"
        f"FIRST 5 SAMPLE ROWS:\n"
    )
    for row in rows[:6]:
        summary += ", ".join(row) + "\n"

    chunks = [{
        "source_type": "dataset_summary",
        "source_location": 1,
        "content": summary
    }]
    return chunks, 1

def parse_excel_dataset(file_id: str, file_bytes: bytes):
    try:
        import pandas as pd
        excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
        chunks = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            summary = (
                f"EXCEL SHEET SUMMARY: '{sheet_name}'\n"
                f"Shape: {df.shape[0]} rows x {df.shape[1]} columns\n"
                f"Columns: {', '.join(map(str, df.columns[:15]))}\n"
                f"Data Types: {df.dtypes.to_dict()}\n"
                f"Null Counts: {df.isnull().sum().to_dict()}\n\n"
                f"SAMPLE DATA:\n{df.head(5).to_string()}\n"
            )
            chunks.append({
                "source_type": "sheet",
                "source_location": sheet_name,
                "content": summary
            })
        return chunks, len(excel_file.sheet_names)
    except Exception as e:
        logger.warning(f"pandas excel parse fallback: {e}")
        return parse_text_document(file_id, file_bytes)

def parse_json_document(file_id: str, file_bytes: bytes):
    text = file_bytes.decode('utf-8', errors='ignore')
    try:
        data = json.loads(text)
        formatted = json.dumps(data, indent=2)
        chunks = [{
            "source_type": "json_structure",
            "source_location": 1,
            "content": formatted[:3000]
        }]
        return chunks, 1
    except Exception:
        return parse_text_document(file_id, file_bytes)

def parse_text_document(file_id: str, file_bytes: bytes):
    text = file_bytes.decode('utf-8', errors='ignore')
    lines = text.splitlines()
    chunks = []
    chunk_size = 50
    for i in range(0, len(lines), chunk_size):
        chunk_lines = lines[i:i + chunk_size]
        chunks.append({
            "source_type": "section",
            "source_location": (i // chunk_size) + 1,
            "content": "\n".join(chunk_lines)
        })
    return chunks, len(chunks) or 1
